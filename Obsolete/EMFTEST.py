# import pyemf #Doesn't work out the box. LAME
#
# width=8
# height=6
# dpi=300
#
# emf=pyemf.EMF(width,height,dpi)
# thin=emf.CreatePen(pyemf.PS_SOLID,1,(0x01,0x02,0x03))
# emf.SelectObject(thin)
# emf.Polyline([(0,0),(width*dpi,height*dpi)])
# emf.Polyline([(0,height*dpi),(width*dpi,0)])
# emf.save("Export\\test-1.emf")

from pptx import Presentation
from matplotlib import pyplot as plt
import numpy as np
import mplppt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('Export\\test.pptx')

# plot [Line2D]
x = np.linspace(-1,5)
y = np.sin(x)
plt.plot(x,y,color='C1')

# pcolormesh
x = np.linspace(0,1, 100)
y = np.linspace(0,1, 100)
X, Y = np.meshgrid(x,y)
Z = X**2 + Y**2
plt.pcolormesh(X,Y,Z)

# text
text = plt.text(0,0,'hello')

# set limits
plt.ylim(-0.5,1)

# Save figure to pptx
mplppt.savefig('EXPORT\\first_example.pptx')
# show figure
plt.show()